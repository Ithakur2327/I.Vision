import os
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation


def extract_text(file_path: str, file_type: str) -> str:
    file_type = file_type.lower()

    if file_type == "pdf":
        reader = PdfReader(file_path)
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    if file_type == "docx":
        doc = DocxDocument(file_path)
        return "\n".join(p.text for p in doc.paragraphs)

    if file_type == "pptx":
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)

    if file_type in ("txt", "md"):
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    raise ValueError(f"Unsupported file type: {file_type}")


ALLOWED_TYPES = {"pdf", "docx", "pptx", "txt", "md"}
MAX_FILE_SIZE_MB = 25


def validate_file(filename: str, size_bytes: int) -> str:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext not in ALLOWED_TYPES:
        raise ValueError(f"File type '.{ext}' is not supported")
    if size_bytes > MAX_FILE_SIZE_MB * 1024 * 1024:
        raise ValueError(f"File exceeds {MAX_FILE_SIZE_MB}MB limit")
    return ext
